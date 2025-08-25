"""
Python-PPTX Google風デザインテンプレート
Version: 1.0 (Universal Python Design - Final)
Author: Pythonスライド自動生成マスター
Description: 指定されたslide_data配列を元に、Google風デザインに準拠したPowerPointスライドを生成します。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import re

# --- 1. 実行設定 ---
SETTINGS = {
    'OUTPUT_FILE': 'generated_presentation.pptx',
    'TEMPLATE_FILE': None  # Noneの場合は新規作成
}

# --- 2. マスターデザイン設定 (Google Design Ver.) ---
CONFIG = {
    'SLIDE_WIDTH': Inches(10),
    'SLIDE_HEIGHT': Inches(5.625),
    
    'FONTS': {
        'family': 'Arial',
        'sizes': {
            'title': 45,
            'date': 16,
            'section_title': 38,
            'content_title': 28,
            'subhead': 18,
            'body': 14,
            'footer': 9,
            'small': 10,
            'ghost_num': 180
        }
    },
    
    'COLORS': {
        'primary_blue': RGBColor(66, 133, 244),      # #4285F4
        'google_red': RGBColor(234, 67, 53),         # #EA4335
        'google_yellow': RGBColor(251, 188, 4),      # #FBBC04
        'google_green': RGBColor(52, 168, 83),       # #34A853
        'text_primary': RGBColor(51, 51, 51),        # #333333
        'background_white': RGBColor(255, 255, 255), # #FFFFFF
        'background_gray': RGBColor(248, 249, 250),  # #f8f9fa
        'faint_gray': RGBColor(232, 234, 237),       # #e8eaed
        'neutral_gray': RGBColor(158, 158, 158),     # #9e9e9e
        'ghost_gray': RGBColor(239, 239, 237)        # #efefed
    },
    
    'LAYOUTS': {
        'title': {
            'logo': {'left': Inches(0.6), 'top': Inches(1.2), 'width': Inches(1.4)},
            'title': {'left': Inches(0.5), 'top': Inches(2.5), 'width': Inches(8), 'height': Inches(1)},
            'date': {'left': Inches(0.5), 'top': Inches(3.7), 'width': Inches(3), 'height': Inches(0.4)}
        },
        'content': {
            'logo': {'left': Inches(8.5), 'top': Inches(0.2), 'width': Inches(1)},
            'title': {'left': Inches(0.3), 'top': Inches(0.7), 'width': Inches(8), 'height': Inches(0.7)},
            'underline': {'left': Inches(0.3), 'top': Inches(1.4), 'width': Inches(2.7), 'height': Inches(0.05)},
            'subhead': {'left': Inches(0.3), 'top': Inches(1.5), 'width': Inches(8), 'height': Inches(0.3)},
            'body': {'left': Inches(0.3), 'top': Inches(1.9), 'width': Inches(9), 'height': Inches(3.2)}
        },
        'section': {
            'title': {'left': Inches(0.6), 'top': Inches(2.5), 'width': Inches(8.5), 'height': Inches(0.8)},
            'ghost_num': {'left': Inches(0.4), 'top': Inches(1.3), 'width': Inches(3), 'height': Inches(2)}
        }
    },
    
    'FOOTER_TEXT': f'© {2025} Your Organization'
}

# --- 3. スライドデータ（サンプル：必ず置換してください） ---
slide_data = [
    {
        'type': 'title', 
        'title': 'サンプルプレゼンテーション', 
        'date': '2025.08.12', 
        'notes': '本日はお集まりいただきありがとうございます。このプレゼンテーションは、Python-pptx風デザインテンプレートの機能と可能性についてご説明するものです。'
    },
    {
        'type': 'section', 
        'title': '1. はじめに', 
        'notes': '最初のセクションでは、このテンプレートが持つ主要な表現パターンについて概観します。'
    },
    {
        'type': 'cards', 
        'title': 'Python-pptx風デザインのテスト', 
        'subhead': 'モダンなデザインパターン', 
        'columns': 3, 
        'items': [
            {'title': 'パターン1', 'desc': '現状：[[重要機能]]実装済み\n課題：パフォーマンス**最適化**が必要'},
            {'title': 'パターン2', 'desc': '現状：デザイン更新完了\n課題：[[ユーザビリティ改善]]を検討'},
            {'title': 'パターン3', 'desc': '現状：テスト環境構築\n課題：**本番環境への移行準備**'}
        ], 
        'notes': 'こちらがカード形式のスライドです。3つの異なる項目を並べて比較検討する際に便利です。それぞれのカードにはタイトルと説明を設定できます。'
    },
    {
        'type': 'closing', 
        'notes': '以上で説明を終わります。ご清聴ありがとうございました。何かご質問はありますでしょうか。'
    }
]

# --- 4. メイン実行関数 ---
section_counter = 0  # 章番号カウンタ（ゴースト数字用）

def generate_presentation():
    """メインのプレゼンテーション生成関数"""
    global section_counter
    section_counter = 0
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    prs.slide_width = CONFIG['SLIDE_WIDTH']
    prs.slide_height = CONFIG['SLIDE_HEIGHT']
    
    page_counter = 0
    
    for data in slide_data:
        slide_type = data.get('type')
        if slide_type in slide_generators:
            if slide_type not in ['title', 'closing']:
                page_counter += 1
            
            # 空のスライドレイアウトを使用
            slide_layout = prs.slide_layouts[6]  # 空のレイアウト
            slide = prs.slides.add_slide(slide_layout)
            
            # スライド生成
            slide_generators[slide_type](slide, data, page_counter)
            
            # スピーカーノートの設定
            if data.get('notes'):
                notes_slide = slide.notes_slide
                notes_shape = notes_slide.notes_text_frame
                notes_shape.text = data['notes']
    
    # ファイル保存
    prs.save(SETTINGS['OUTPUT_FILE'])
    print(f"プレゼンテーションが生成されました: {SETTINGS['OUTPUT_FILE']}")

# --- 5. スライド生成ディスパッチャ（関数定義の後で設定） ---
slide_generators = {}  # 後で設定

# --- 6. スライド生成関数群 ---
def create_title_slide(slide, data, page_num):
    """タイトルスライドの生成"""
    set_slide_background(slide, 'background_white')
    draw_logo(slide, CONFIG['LAYOUTS']['title']['logo'], size=14)
    
    # タイトル
    title_layout = CONFIG['LAYOUTS']['title']['title']
    title_shape = slide.shapes.add_textbox(
        title_layout['left'], title_layout['top'],
        title_layout['width'], title_layout['height']
    )
    title_frame = title_shape.text_frame
    title_frame.text = data.get('title', '')
    format_text(title_frame, size=CONFIG['FONTS']['sizes']['title'], bold=True)
    
    # 日付
    date_layout = CONFIG['LAYOUTS']['title']['date']
    date_shape = slide.shapes.add_textbox(
        date_layout['left'], date_layout['top'],
        date_layout['width'], date_layout['height']
    )
    date_frame = date_shape.text_frame
    date_frame.text = data.get('date', '')
    format_text(date_frame, size=CONFIG['FONTS']['sizes']['date'])
    
    # 下部バー
    draw_bottom_bar(slide)

def create_section_slide(slide, data, page_num):
    """章扉スライドの生成"""
    global section_counter
    
    set_slide_background(slide, 'background_gray')
    
    section_counter += 1
    
    # ゴースト番号の生成
    section_no = data.get('section_no', section_counter)
    if 'section_no' not in data:
        # タイトルから番号を抽出
        title = data.get('title', '')
        match = re.match(r'^\s*(\d+)[\.．]', title)
        if match:
            section_no = int(match.group(1))
    
    ghost_num = f"{section_no:02d}"
    
    # ゴースト番号
    ghost_layout = CONFIG['LAYOUTS']['section']['ghost_num']
    ghost_shape = slide.shapes.add_textbox(
        ghost_layout['left'], ghost_layout['top'],
        ghost_layout['width'], ghost_layout['height']
    )
    ghost_frame = ghost_shape.text_frame
    ghost_frame.text = ghost_num
    format_text(ghost_frame, 
                size=CONFIG['FONTS']['sizes']['ghost_num'], 
                bold=True, 
                color=CONFIG['COLORS']['ghost_gray'])
    
    # タイトル
    title_layout = CONFIG['LAYOUTS']['section']['title']
    title_shape = slide.shapes.add_textbox(
        title_layout['left'], title_layout['top'],
        title_layout['width'], title_layout['height']
    )
    title_frame = title_shape.text_frame
    title_frame.text = data.get('title', '')
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    format_text(title_frame, size=CONFIG['FONTS']['sizes']['section_title'], bold=True)
    
    # フッター
    add_footer(slide, page_num)

def create_content_slide(slide, data, page_num):
    """コンテンツスライドの生成"""
    set_slide_background(slide, 'background_white')
    
    # 標準ヘッダー
    draw_standard_header(slide, data.get('title', ''))
    
    # サブヘッダー
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # アジェンダ安全装置
    points = data.get('points', [])
    if is_agenda_title(data.get('title', '')) and not points:
        points = build_agenda_from_slide_data() or ['本日の目的', '進め方', '次のアクション']
    
    # コンテンツ描画
    if points:
        body_layout = CONFIG['LAYOUTS']['content']['body']
        body_top = body_layout['top'] + Inches(subhead_offset)
        
        if data.get('two_column') or data.get('columns'):
            # 2カラム表示
            draw_two_column_content(slide, points, data.get('columns'), body_top)
        else:
            # 1カラム表示
            draw_single_column_content(slide, points, body_top)
    
    # 画像表示
    if data.get('images'):
        draw_images(slide, data['images'])
    
    # フッターとボトムバー
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_compare_slide(slide, data, page_num):
    """比較スライドの生成"""
    set_slide_background(slide, 'background_white')
    
    # 標準ヘッダー
    draw_standard_header(slide, data.get('title', ''))
    
    # サブヘッダー
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # 左右比較ボックス
    body_layout = CONFIG['LAYOUTS']['content']['body']
    body_top = body_layout['top'] + Inches(subhead_offset)
    
    # 左ボックス
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), body_top,
        Inches(4.2), Inches(3.2)
    )
    format_compare_box(left_box, data.get('left_title', '選択肢A'), data.get('left_items', []))
    
    # 右ボックス
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.2), body_top,
        Inches(4.2), Inches(3.2)
    )
    format_compare_box(right_box, data.get('right_title', '選択肢B'), data.get('right_items', []))
    
    # フッターとボトムバー
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_process_slide(slide, data, page_num):
    """プロセススライドの生成"""
    set_slide_background(slide, 'background_white')
    
    draw_standard_header(slide, data.get('title', ''))
    
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # プロセスステップの描画
    steps = data.get('steps', [])
    if steps:
        draw_process_steps(slide, steps, subhead_offset)
    
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_timeline_slide(slide, data, page_num):
    """タイムラインスライドの生成"""
    set_slide_background(slide, 'background_white')
    
    draw_standard_header(slide, data.get('title', ''))
    
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # タイムライン描画
    milestones = data.get('milestones', [])
    if milestones:
        draw_timeline(slide, milestones, subhead_offset)
    
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_diagram_slide(slide, data, page_num):
    """ダイアグラムスライドの生成"""
    set_slide_background(slide, 'background_white')
    
    draw_standard_header(slide, data.get('title', ''))
    
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # レーン図の描画
    lanes = data.get('lanes', [])
    if lanes:
        draw_diagram_lanes(slide, lanes, subhead_offset)
    
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_cards_slide(slide, data, page_num):
    """カードスライドの生成"""
    set_slide_background(slide, 'background_white')
    
    draw_standard_header(slide, data.get('title', ''))
    
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # カードグリッドの描画
    items = data.get('items', [])
    columns = data.get('columns', 2 if len(items) <= 4 else 3)
    if items:
        draw_cards_grid(slide, items, columns, subhead_offset)
    
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_table_slide(slide, data, page_num):
    """テーブルスライドの生成"""
    set_slide_background(slide, 'background_white')
    
    draw_standard_header(slide, data.get('title', ''))
    
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # テーブルの描画
    headers = data.get('headers', [])
    rows = data.get('rows', [])
    if headers:
        draw_table(slide, headers, rows, subhead_offset)
    
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_progress_slide(slide, data, page_num):
    """進捗スライドの生成"""
    set_slide_background(slide, 'background_white')
    
    draw_standard_header(slide, data.get('title', ''))
    
    subhead_offset = 0
    if data.get('subhead'):
        subhead_offset = draw_subhead(slide, data['subhead'])
    
    # 進捗バーの描画
    items = data.get('items', [])
    if items:
        draw_progress_bars(slide, items, subhead_offset)
    
    draw_bottom_bar(slide)
    add_footer(slide, page_num)

def create_closing_slide(slide, data, page_num):
    """クロージングスライドの生成"""
    set_slide_background(slide, 'background_white')
def set_slide_background(slide, color_key):
    """背景色を設定する共通関数"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CONFIG['COLORS'][color_key]

def draw_logo(slide, layout, size=12):
    """ロゴ描画の共通関数"""
    logo_shape = slide.shapes.add_textbox(
        layout['left'], layout['top'],
        layout['width'], Inches(0.5)
    )
    logo_frame = logo_shape.text_frame
    logo_frame.text = "LOGO"
    format_text(logo_frame, size=size, bold=True, color=CONFIG['COLORS']['primary_blue'])
    
    # 中央にロゴ（テキスト）を配置
    logo_shape = slide.shapes.add_textbox(
        Inches(3.5), Inches(2.5),
        Inches(3), Inches(1)
    )
    logo_frame = logo_shape.text_frame
    logo_frame.text = "THANK YOU"
    logo_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    format_text(logo_frame, size=36, bold=True, color=CONFIG['COLORS']['primary_blue'])

# --- 7. ユーティリティ関数群 ---
def format_text(text_frame, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    """テキストのフォーマット設定"""
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    
    run = paragraph.runs[0]
    font = run.font
    font.name = CONFIG['FONTS']['family']
    font.size = Pt(size)
    font.bold = bold
    
    if color:
        font.color.rgb = color
    else:
        font.color.rgb = CONFIG['COLORS']['text_primary']

def apply_inline_styles(text_frame, text):
    """インライン装飾の適用"""
    # **太字** パターン
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    
    # [[重要語]] パターン
    text = re.sub(r'\[\[(.*?)\]\]', r'\1', text)
    
    text_frame.text = text
    
    # スタイル適用は簡易版（実際の実装では範囲指定が必要）
    return text

def draw_standard_header(slide, title):
    """標準ヘッダーの描画"""
    # ロゴ
    logo_layout = CONFIG['LAYOUTS']['content']['logo']
    logo_shape = slide.shapes.add_textbox(
        logo_layout['left'], logo_layout['top'],
        logo_layout['width'], Inches(0.3)
    )
    logo_frame = logo_shape.text_frame
    logo_frame.text = "LOGO"
    format_text(logo_frame, size=12, color=CONFIG['COLORS']['primary_blue'])
    
    # タイトル
    title_layout = CONFIG['LAYOUTS']['content']['title']
    title_shape = slide.shapes.add_textbox(
        title_layout['left'], title_layout['top'],
        title_layout['width'], title_layout['height']
    )
    title_frame = title_shape.text_frame
    title_frame.text = apply_inline_styles(title_frame, title)
    format_text(title_frame, size=CONFIG['FONTS']['sizes']['content_title'], bold=True)
    
    # アンダーライン
    underline_layout = CONFIG['LAYOUTS']['content']['underline']
    underline_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        underline_layout['left'], underline_layout['top'],
        underline_layout['width'], underline_layout['height']
    )
    fill = underline_shape.fill
    fill.solid()
    fill.fore_color.rgb = CONFIG['COLORS']['primary_blue']
    underline_shape.line.fill.background()

def draw_subhead(slide, subhead):
    """サブヘッドの描画"""
    subhead_layout = CONFIG['LAYOUTS']['content']['subhead']
    subhead_shape = slide.shapes.add_textbox(
        subhead_layout['left'], subhead_layout['top'],
        subhead_layout['width'], subhead_layout['height']
    )
    subhead_frame = subhead_shape.text_frame
    subhead_frame.text = apply_inline_styles(subhead_frame, subhead)
    format_text(subhead_frame, size=CONFIG['FONTS']['sizes']['subhead'])
    return 0.4  # オフセット値

def draw_single_column_content(slide, points, top_offset):
    """単一カラムコンテンツの描画"""
    body_layout = CONFIG['LAYOUTS']['content']['body']
    content_shape = slide.shapes.add_textbox(
        body_layout['left'], top_offset,
        body_layout['width'], body_layout['height']
    )
    content_frame = content_shape.text_frame
    
    for i, point in enumerate(points):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = f"• {apply_inline_styles(content_frame, point)}"
        format_text_paragraph(p, CONFIG['FONTS']['sizes']['body'])

def draw_two_column_content(slide, points, columns, top_offset):
    """2カラムコンテンツの描画"""
    if columns and len(columns) == 2:
        left_items, right_items = columns[0], columns[1]
    else:
        mid = len(points) // 2
        left_items, right_items = points[:mid], points[mid:]
    
    # 左カラム
    left_shape = slide.shapes.add_textbox(
        Inches(0.3), top_offset,
        Inches(4.2), Inches(3.2)
    )
    draw_bullet_points(left_shape.text_frame, left_items)
    
    # 右カラム
    right_shape = slide.shapes.add_textbox(
        Inches(5.2), top_offset,
        Inches(4.2), Inches(3.2)
    )
    draw_bullet_points(right_shape.text_frame, right_items)

def draw_bullet_points(text_frame, points):
    """箇条書きポイントの描画"""
    for i, point in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = f"• {point}"
        format_text_paragraph(p, CONFIG['FONTS']['sizes']['body'])

def format_text_paragraph(paragraph, size):
    """段落のテキストフォーマット"""
    run = paragraph.runs[0] if paragraph.runs else paragraph.runs.add()
    font = run.font
    font.name = CONFIG['FONTS']['family']
    font.size = Pt(size)
    font.color.rgb = CONFIG['COLORS']['text_primary']

def format_compare_box(shape, title, items):
    """比較ボックスのフォーマット"""
    # 背景色設定
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CONFIG['COLORS']['background_gray']
    
    # 枠線設定
    line = shape.line
    line.color.rgb = CONFIG['COLORS']['faint_gray']
    line.width = Pt(1)
    
    # テキスト設定
    text_frame = shape.text_frame
    text_frame.text = title
    if items:
        text_frame.text += "\n\n" + "\n".join([f"• {item}" for item in items])
    
    format_text(text_frame, size=CONFIG['FONTS']['sizes']['body'], bold=True)

def draw_process_steps(slide, steps, subhead_offset):
    """プロセスステップの描画"""
    body_layout = CONFIG['LAYOUTS']['content']['body']
    start_top = body_layout['top'] + Inches(subhead_offset)
    
    step_height = Inches(0.6)
    for i, step in enumerate(steps):
        y_pos = start_top + i * step_height
        
        # ステップ番号
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), y_pos,
            Inches(0.4), Inches(0.4)
        )
        fill = num_shape.fill
        fill.solid()
        fill.fore_color.rgb = CONFIG['COLORS']['primary_blue']
        
        num_frame = num_shape.text_frame
        num_frame.text = str(i + 1)
        format_text(num_frame, size=12, bold=True, color=CONFIG['COLORS']['background_white'], align=PP_ALIGN.CENTER)
        
        # ステップテキスト
        text_shape = slide.shapes.add_textbox(
            Inches(1), y_pos,
            Inches(8), Inches(0.4)
        )
        text_frame = text_shape.text_frame
        text_frame.text = apply_inline_styles(text_frame, step)
        format_text(text_frame, size=CONFIG['FONTS']['sizes']['body'])

def draw_timeline(slide, milestones, subhead_offset):
    """タイムラインの描画"""
    if not milestones:
        return
    
    # タイムラインの基準位置
    timeline_y = Inches(3.2)
    timeline_left = Inches(1)
    timeline_width = Inches(8)
    
    # ベースライン
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        timeline_left, timeline_y,
        timeline_width, Inches(0.05)
    )
    fill = line_shape.fill
    fill.solid()
    fill.fore_color.rgb = CONFIG['COLORS']['faint_gray']
    
    # マイルストーン
    gap = timeline_width / max(1, len(milestones) - 1) if len(milestones) > 1 else Inches(0)
    
    for i, milestone in enumerate(milestones):
        x_pos = timeline_left + i * gap if len(milestones) > 1 else timeline_left + timeline_width / 2
        
        # ドット
        dot_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x_pos - Inches(0.1), timeline_y - Inches(0.1),
            Inches(0.2), Inches(0.2)
        )
        
        state = milestone.get('state', 'todo').lower()
        dot_fill = dot_shape.fill
        dot_fill.solid()
        
        if state == 'done':
            dot_fill.fore_color.rgb = CONFIG['COLORS']['google_green']
        elif state == 'next':
            dot_fill.fore_color.rgb = CONFIG['COLORS']['google_yellow']
        else:
            dot_fill.fore_color.rgb = CONFIG['COLORS']['neutral_gray']
        
        # ラベル
        label_shape = slide.shapes.add_textbox(
            x_pos - Inches(0.8), timeline_y - Inches(0.6),
            Inches(1.6), Inches(0.3)
        )
        label_frame = label_shape.text_frame
        label_frame.text = milestone.get('label', '')
        format_text(label_frame, size=CONFIG['FONTS']['sizes']['small'], align=PP_ALIGN.CENTER)
        
        # 日付
        date_shape = slide.shapes.add_textbox(
            x_pos - Inches(0.8), timeline_y + Inches(0.3),
            Inches(1.6), Inches(0.3)
        )
        date_frame = date_shape.text_frame
        date_frame.text = milestone.get('date', '')
        format_text(date_frame, size=CONFIG['FONTS']['sizes']['small'], 
                   color=CONFIG['COLORS']['neutral_gray'], align=PP_ALIGN.CENTER)

def draw_diagram_lanes(slide, lanes, subhead_offset):
    """レーン図の描画"""
    if not lanes:
        return
    
    body_layout = CONFIG['LAYOUTS']['content']['body']
    start_top = body_layout['top'] + Inches(subhead_offset)
    
    lane_width = Inches(8.5) / len(lanes)
    
    for i, lane in enumerate(lanes):
        x_pos = Inches(0.5) + i * lane_width
        
        # レーンタイトル
        title_shape = slide.shapes.add_textbox(
            x_pos, start_top,
            lane_width - Inches(0.1), Inches(0.4)
        )
        title_frame = title_shape.text_frame
        title_frame.text = lane.get('title', '')
        format_text(title_frame, size=CONFIG['FONTS']['sizes']['body'], bold=True, align=PP_ALIGN.CENTER)
        
        # アイテム
        items = lane.get('items', [])
        item_height = Inches(0.4)
        
        for j, item in enumerate(items):
            item_y = start_top + Inches(0.6) + j * item_height
            item_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos + Inches(0.05), item_y,
                lane_width - Inches(0.2), item_height - Inches(0.05)
            )
            
            # カードスタイリング
            fill = item_shape.fill
            fill.solid()
            fill.fore_color.rgb = CONFIG['COLORS']['background_white']
            
            line = item_shape.line
            line.color.rgb = CONFIG['COLORS']['faint_gray']
            line.width = Pt(1)
            
            item_frame = item_shape.text_frame
            item_frame.text = item
            format_text(item_frame, size=CONFIG['FONTS']['sizes']['small'], align=PP_ALIGN.CENTER)

def draw_cards_grid(slide, items, columns, subhead_offset):
    """カードグリッドの描画"""
    if not items:
        return
    
    body_layout = CONFIG['LAYOUTS']['content']['body']
    start_top = body_layout['top'] + Inches(subhead_offset)
    
    rows = (len(items) + columns - 1) // columns
    card_width = Inches(8.5) / columns
    card_height = Inches(2.8) / rows
    
    for i, item in enumerate(items):
        row = i // columns
        col = i % columns
        
        x_pos = Inches(0.5) + col * card_width
        y_pos = start_top + row * card_height
        
        # カード
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos + Inches(0.05), y_pos + Inches(0.05),
            card_width - Inches(0.1), card_height - Inches(0.1)
        )
        
        # カードスタイリング
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = CONFIG['COLORS']['background_white']
        
        line = card_shape.line
        line.color.rgb = CONFIG['COLORS']['faint_gray']
        line.width = Pt(1)
        
        # テキスト
        card_frame = card_shape.text_frame
        
        if isinstance(item, dict):
            title = item.get('title', '')
            desc = item.get('desc', '')
            card_frame.text = f"{title}\n{desc}" if desc else title
        else:
            card_frame.text = str(item)
        
        card_frame.text = apply_inline_styles(card_frame, card_frame.text)
        format_text(card_frame, size=CONFIG['FONTS']['sizes']['small'])

def draw_table(slide, headers, rows, subhead_offset):
    """テーブルの描画"""
    if not headers:
        return
    
    body_layout = CONFIG['LAYOUTS']['content']['body']
    start_top = body_layout['top'] + Inches(subhead_offset)
    
    col_width = Inches(8.5) / len(headers)
    row_height = Inches(0.4)
    
    # ヘッダー
    for i, header in enumerate(headers):
        x_pos = Inches(0.5) + i * col_width
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x_pos, start_top,
            col_width, row_height
        )
        
        fill = header_shape.fill
        fill.solid()
        fill.fore_color.rgb = CONFIG['COLORS']['primary_blue']
        
        header_frame = header_shape.text_frame
        header_frame.text = str(header)
        format_text(header_frame, size=CONFIG['FONTS']['sizes']['small'], 
                   bold=True, color=CONFIG['COLORS']['background_white'], align=PP_ALIGN.CENTER)
    
    # データ行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            x_pos = Inches(0.5) + col_idx * col_width
            y_pos = start_top + (row_idx + 1) * row_height
            
            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x_pos, y_pos,
                col_width, row_height
            )
            
            fill = cell_shape.fill
            fill.solid()
            fill.fore_color.rgb = CONFIG['COLORS']['background_white']
            
            line = cell_shape.line
            line.color.rgb = CONFIG['COLORS']['faint_gray']
            line.width = Pt(1)
            
            cell_frame = cell_shape.text_frame
            cell_frame.text = str(cell_data) if cell_data else ''
            format_text(cell_frame, size=CONFIG['FONTS']['sizes']['small'], align=PP_ALIGN.CENTER)

def draw_progress_bars(slide, items, subhead_offset):
    """進捗バーの描画"""
    if not items:
        return
    
    body_layout = CONFIG['LAYOUTS']['content']['body']
    start_top = body_layout['top'] + Inches(subhead_offset)
    
    bar_height = Inches(0.3)
    row_height = Inches(0.6)
    
    for i, item in enumerate(items):
        y_pos = start_top + i * row_height
        
        # ラベル
        label_shape = slide.shapes.add_textbox(
            Inches(0.5), y_pos,
            Inches(2.5), bar_height
        )
        label_frame = label_shape.text_frame
        label_frame.text = item.get('label', '')
        format_text(label_frame, size=CONFIG['FONTS']['sizes']['body'])
        
        # 背景バー
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.2), y_pos + Inches(0.05),
            Inches(5), Inches(0.2)
        )
        bg_fill = bg_bar.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = CONFIG['COLORS']['faint_gray']
        bg_bar.line.fill.background()
        
        # 進捗バー
        percent = max(0, min(100, item.get('percent', 0)))
        progress_width = Inches(5) * (percent / 100)
        
        if progress_width > 0:
            progress_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3.2), y_pos + Inches(0.05),
                progress_width, Inches(0.2)
            )
            progress_fill = progress_bar.fill
            progress_fill.solid()
            progress_fill.fore_color.rgb = CONFIG['COLORS']['google_green']
            progress_bar.line.fill.background()
        
        # パーセンテージ
        pct_shape = slide.shapes.add_textbox(
            Inches(8.5), y_pos,
            Inches(1), bar_height
        )
        pct_frame = pct_shape.text_frame
        pct_frame.text = f"{percent}%"
        format_text(pct_frame, size=CONFIG['FONTS']['sizes']['small'], 
                   color=CONFIG['COLORS']['neutral_gray'])

def draw_images(slide, images):
    """画像の描画"""
    # 簡易的な画像表示（実装はURLからの画像取得が必要）
    pass

def draw_bottom_bar(slide):
    """下部バーの描画"""
    bar_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), CONFIG['SLIDE_HEIGHT'] - Inches(0.1),
        CONFIG['SLIDE_WIDTH'], Inches(0.1)
    )
    fill = bar_shape.fill
    fill.solid()
    fill.fore_color.rgb = CONFIG['COLORS']['primary_blue']
    bar_shape.line.fill.background()

def add_footer(slide, page_num):
    """フッターの追加"""
    # 左側テキスト
    left_footer = slide.shapes.add_textbox(
        Inches(0.2), CONFIG['SLIDE_HEIGHT'] - Inches(0.3),
        Inches(4), Inches(0.2)
    )
    left_frame = left_footer.text_frame
    left_frame.text = CONFIG['FOOTER_TEXT']
    format_text(left_frame, size=CONFIG['FONTS']['sizes']['footer'])
    
    # ページ番号（右側）
    if page_num > 0:
        right_footer = slide.shapes.add_textbox(
            CONFIG['SLIDE_WIDTH'] - Inches(1), CONFIG['SLIDE_HEIGHT'] - Inches(0.3),
            Inches(0.8), Inches(0.2)
        )
        right_frame = right_footer.text_frame
        right_frame.text = str(page_num)
        format_text(right_frame, size=CONFIG['FONTS']['sizes']['footer'], 
                    color=CONFIG['COLORS']['primary_blue'], align=PP_ALIGN.RIGHT)

# ユーティリティ関数
def is_agenda_title(title):
    """アジェンダタイトルの判定"""
    t = title.lower()
    return any(keyword in t for keyword in ['agenda', 'アジェンダ', '目次', '本日お伝えすること'])

def build_agenda_from_slide_data():
    """slide_dataからアジェンダを構築"""
    agenda_items = []
    for data in slide_data:
        if data.get('type') == 'section' and data.get('title'):
            agenda_items.append(data['title'].strip())
    
    return agenda_items[:5] if agenda_items else []

# slide_generators辞書を関数定義後に設定
slide_generators = {
    'title': create_title_slide,
    'section': create_section_slide,
    'content': create_content_slide,
    'compare': create_compare_slide,
    'process': create_process_slide,
    'timeline': create_timeline_slide,
    'diagram': create_diagram_slide,
    'cards': create_cards_slide,
    'table': create_table_slide,
    'progress': create_progress_slide,
    'closing': create_closing_slide
}

# --- 8. メイン実行 ---
if __name__ == '__main__':
    generate_presentation()