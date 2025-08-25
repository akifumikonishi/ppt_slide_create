from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

slide_data = [
    {
        'type': 'title',
        'title': 'SuccessFactors グローバルロールアウト\n日本・南アフリカ統合プロジェクト リスク分析',
        'date': '2025.08.25',
        'notes': 'SuccessFactorsのグローバルロールアウトにおける日本本社と南アフリカ子会社の統合プロジェクトについて、2024-2025年の実装事例から抽出したリスクと対処策をご紹介します。本調査では、10件以上の導入事例を分析し、共通する重要要素をピックアップしました。'
    },
    {
        'type': 'content',
        'title': '調査概要：2024-2025年 SuccessFactors 導入事例分析',
        'subhead': '10件以上の実装事例から抽出した共通リスクと成功要因',
        'points': [
            '**分析対象：** 2024年から2025年のSuccessFactors導入プロジェクト',
            '**調査範囲：** グローバル企業のマルチインスタンス統合事例',
            '**特定リスク：** 日本本社・南アフリカ子会社統合における課題',
            '**抽出項目：** 技術的リスク、運用リスク、組織的リスク',
            '**成功要因：** プロジェクト成功につながる対処策とベストプラクティス'
        ],
        'notes': '本調査では、2024年から2025年にかけてのSuccessFactors導入事例を詳細に分析しました。特に、日本本社で使用中のSuccessFactorsに南アフリカ子会社のシステムを統合する際の特有のリスクに焦点を当てています。グローバル企業200社以上の法人を対象とした実装事例や、アフリカ最大級の保険会社での導入事例などを参考にしています。'
    },
    {
        'type': 'compare',
        'title': '主要リスクカテゴリー分析',
        'subhead': '技術的リスクと運用リスクの対比',
        'left_title': '技術的リスク',
        'right_title': '運用リスク',
        'left_items': [
            '**マルチインスタンス統合**\n複数のEmployee Central統合時の技術課題',
            '**認証システム移行**\n2025年6月OAuth 2.0必須対応',
            '**APIデプリケーション**\nSOAP API廃止への対応（2025年5月）',
            '**データ整合性確保**\nシステム間データ同期の複雑性',
            '**レポーティング統合**\n異なる地域間での統一レポート'
        ],
        'right_items': [
            '**変革管理**\nユーザー受け入れとトレーニング課題',
            '**段階的ロールアウト**\n移行期間中の業務継続リスク',
            '**地域固有要件**\n日本・南アフリカの法的規制対応',
            '**データ品質管理**\nレガシーデータの変換・検証',
            '**プロジェクト体制**\n専門スキル不足とパートナー依存'
        ],
        'notes': '技術的リスクでは、2025年に予定されている大規模なシステム変更への対応が急務となっています。特に認証システムの移行とAPI変更は必須対応項目です。運用リスクでは、グローバル企業特有の地域差への対応と、段階的ロールアウト時の業務継続性確保が重要となります。'
    },
    {
        'type': 'timeline',
        'title': '2025年 必須対応スケジュール',
        'subhead': 'SAP SuccessFactorsシステム変更に伴う対応期限',
        'milestones': [
            {'label': 'SOAP API廃止対応', 'date': '2025年5月', 'state': 'next'},
            {'label': 'SSO証明書期限', 'date': '2025年6月', 'state': 'next'},
            {'label': '基本認証サポート終了', 'date': '2025年6月', 'state': 'next'},
            {'label': 'レポート機能移行', 'date': '2026年後半', 'state': 'todo'},
            {'label': '基本認証完全削除', 'date': '2026年11月', 'state': 'todo'}
        ],
        'notes': '2025年は SuccessFactors にとって大きな変革の年となります。5月にSOAP APIが廃止され、6月にはシングルサインオン証明書の更新と基本認証のサポート終了が予定されています。これらの変更は必須対応であり、期限内に対応しない場合、システムアクセスができなくなるリスクがあります。プロジェクト計画時にはこれらの期限を必ず考慮する必要があります。'
    },
    {
        'type': 'cards',
        'title': '統合成功のための対処策',
        'subhead': '実装事例から導出された効果的な対策',
        'columns': 2,
        'items': [
            {
                'title': '**段階的移行戦略**',
                'desc': '一括移行ではなく段階的アプローチで\nリスクを分散・最小化'
            },
            {
                'title': '**専用テスト環境**',
                'desc': 'データ移行専用環境を構築し\n安定した検証環境を確保'
            },
            {
                'title': '**データガバナンス強化**',
                'desc': '移行データの品質管理と\n責任体制の明確化'
            },
            {
                'title': '**変革管理プログラム**',
                'desc': 'ユーザートレーニングと\nコミュニケーション戦略'
            },
            {
                'title': '**OAuth 2.0 早期移行**',
                'desc': '2025年期限前の\n認証システム更新実施'
            },
            {
                'title': '**統合パートナー活用**',
                'desc': '経験豊富な実装パートナーとの\n協業体制構築'
            }
        ],
        'notes': '成功事例から抽出された対処策では、特に段階的移行戦略の重要性が強調されています。一括でのシステム統合は失敗リスクが高く、段階的なアプローチにより各段階でのリスク評価と対応が可能になります。また、データ移行の成功率を向上させるため、専用のテスト環境の構築が推奨されています。Gartner調査では、データ移行プロジェクトの83%が失敗または予算・スケジュール超過していることから、これらの対策は必須と考えられます。'
    },
    {
        'type': 'closing',
        'notes': '本調査により、SuccessFactorsグローバルロールアウトにおける主要なリスクと効果的な対処策が明確になりました。特に2025年の必須対応項目への早期着手と、段階的移行戦略の採用が成功の鍵となります。日本本社と南アフリカ子会社の統合プロジェクトにおいては、地域固有の要件に配慮しながら、グローバル標準に準拠したアプローチが重要です。ご質問がございましたら、お気軽にお声かけください。'
    }
]

def create_title_slide(prs, slide_data):
    """タイトルスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空のレイアウト
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 日付
    date_box = slide.shapes.add_textbox(Inches(13), Inches(8), Inches(2.5), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = slide_data['date']
    date_frame.paragraphs[0].font.size = Pt(14)
    date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def create_content_slide(prs, slide_data):
    """コンテンツスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # サブヘッド
    if 'subhead' in slide_data:
        subhead_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
        subhead_frame = subhead_box.text_frame
        subhead_frame.text = slide_data['subhead']
        subhead_frame.paragraphs[0].font.size = Pt(18)
        subhead_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    # ポイント
    if 'points' in slide_data:
        points_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(5.5))
        points_frame = points_box.text_frame
        for i, point in enumerate(slide_data['points']):
            if i > 0:
                p = points_frame.add_paragraph()
            else:
                p = points_frame.paragraphs[0]
            
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)

def create_compare_slide(prs, slide_data):
    """比較スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # サブヘッド
    if 'subhead' in slide_data:
        subhead_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
        subhead_frame = subhead_box.text_frame
        subhead_frame.text = slide_data['subhead']
        subhead_frame.paragraphs[0].font.size = Pt(18)
        subhead_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    # 左側タイトル
    left_title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6.5), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = slide_data['left_title']
    left_title_frame.paragraphs[0].font.size = Pt(22)
    left_title_frame.paragraphs[0].font.bold = True
    left_title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # 右側タイトル
    right_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.5), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = slide_data['right_title']
    right_title_frame.paragraphs[0].font.size = Pt(22)
    right_title_frame.paragraphs[0].font.bold = True
    right_title_frame.paragraphs[0].font.color.rgb = RGBColor(204, 102, 0)
    
    # 左側アイテム
    left_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(6.5), Inches(5))
    left_frame = left_box.text_frame
    for i, item in enumerate(slide_data['left_items']):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # 右側アイテム
    right_box = slide.shapes.add_textbox(Inches(8.5), Inches(3.2), Inches(6.5), Inches(5))
    right_frame = right_box.text_frame
    for i, item in enumerate(slide_data['right_items']):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)

def create_timeline_slide(prs, slide_data):
    """タイムラインスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # サブヘッド
    if 'subhead' in slide_data:
        subhead_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
        subhead_frame = subhead_box.text_frame
        subhead_frame.text = slide_data['subhead']
        subhead_frame.paragraphs[0].font.size = Pt(18)
        subhead_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    # マイルストーン
    y_pos = 3
    for milestone in slide_data['milestones']:
        # 状態に応じた色設定
        if milestone['state'] == 'done':
            color = RGBColor(46, 125, 50)
        elif milestone['state'] == 'next':
            color = RGBColor(255, 152, 0)
        else:  # todo
            color = RGBColor(158, 158, 158)
        
        # 日付
        date_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = milestone['date']
        date_frame.paragraphs[0].font.size = Pt(18)
        date_frame.paragraphs[0].font.bold = True
        date_frame.paragraphs[0].font.color.rgb = color
        
        # ラベル
        label_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(10), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = milestone['label']
        label_frame.paragraphs[0].font.size = Pt(18)
        label_frame.paragraphs[0].font.color.rgb = color
        
        y_pos += 0.8

def create_cards_slide(prs, slide_data):
    """カードスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # サブヘッド
    if 'subhead' in slide_data:
        subhead_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
        subhead_frame = subhead_box.text_frame
        subhead_frame.text = slide_data['subhead']
        subhead_frame.paragraphs[0].font.size = Pt(18)
        subhead_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    # カード配置
    columns = slide_data.get('columns', 2)
    card_width = 14 / columns - 0.5
    card_height = 2
    
    for i, item in enumerate(slide_data['items']):
        col = i % columns
        row = i // columns
        
        x = 1 + col * (card_width + 0.5)
        y = 2.8 + row * (card_height + 0.3)
        
        # カード背景
        card_shape = slide.shapes.add_shape(
            1,  # rectangle
            Inches(x), Inches(y), Inches(card_width), Inches(card_height)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        card_shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # カードテキスト
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.2), 
            Inches(card_width - 0.4), Inches(card_height - 0.4)
        )
        text_frame = text_box.text_frame
        
        if isinstance(item, dict):
            text_frame.text = item['title']
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            if 'desc' in item:
                p = text_frame.add_paragraph()
                p.text = item['desc']
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(102, 102, 102)
        else:
            text_frame.text = item
            text_frame.paragraphs[0].font.size = Pt(14)

def create_closing_slide(prs, slide_data):
    """クロージングスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 「ありがとうございました」
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "ありがとうございました"
    thanks_frame.paragraphs[0].font.size = Pt(48)
    thanks_frame.paragraphs[0].font.bold = True
    thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    thanks_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# スライド生成
for slide_info in slide_data:
    if slide_info['type'] == 'title':
        create_title_slide(prs, slide_info)
    elif slide_info['type'] == 'content':
        create_content_slide(prs, slide_info)
    elif slide_info['type'] == 'compare':
        create_compare_slide(prs, slide_info)
    elif slide_info['type'] == 'timeline':
        create_timeline_slide(prs, slide_info)
    elif slide_info['type'] == 'cards':
        create_cards_slide(prs, slide_info)
    elif slide_info['type'] == 'closing':
        create_closing_slide(prs, slide_info)

# ディレクトリ作成
output_dir = r"C:\Users\022038760\work\ppt_slide_create\01_分析レポート"
os.makedirs(output_dir, exist_ok=True)

# ファイル保存
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
filename = f"SuccessFactorsグローバルロールアウト事例_{timestamp}.pptx"
filepath = os.path.join(output_dir, filename)
prs.save(filepath)

print(f"プレゼンテーションが正常に作成されました：{filepath}")